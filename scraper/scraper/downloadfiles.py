import os
import zipfile
import tempfile
from docx import Document
from oletools.olevba3 import VBA_Parser
import fitz  # PyMuPDF
import subprocess
import shutil
from urllib.parse import urlparse
import re
import pptx
import openpyxl
import csv
from urllib.parse import urlparse, urlunparse


# ---------- Config ----------
TO_PROCESS_FURTHER_PATH = "scraper/scraper/toProcessFurther"
SUPPORTED_EXTENSIONS = {
    ".docx",
    ".doc",
    ".pdf",
    ".zip",
    ".pptx",
    ".ppt",
    ".xlsx",
    ".csv",
}


# Extractors
def extract_links_from_docx(filepath):
    """Extracts hyperlinks and visible URLs from a .docx file."""

    links = []
    try:
        doc = Document(filepath)
        for rel in doc.part.rels.values():
            if "hyperlink" in rel.reltype.lower() and "http" in rel.target_ref:
                link = rel.target_ref.strip()
                links.append(link)
        for para in doc.paragraphs:
            found_urls = re.findall(r"https?://\S+|www\.\S+", para.text)
            for url in found_urls:
                cleaned_url = url.rstrip(".,;)!?\"'")
                links.append(cleaned_url)
    except Exception as e:
        print(f"[WARNING] Error reading .docx file {filepath}: {e}")
    return links


def extract_links_from_doc(filepath):
    """Extracts URLs embedded in VBA macros from a legacy .doc file."""
    links = []
    vba = None
    try:
        vba = VBA_Parser(filepath)
        if vba.detect_vba_macros():
            for _, _, vba_code in vba.extract_macros():
                for line in vba_code.splitlines():
                    if "http" in line or "www." in line:
                        links.append(line.strip())
    except Exception as e:
        print(f"[WARNING] Error reading .doc file {filepath}: {e}")
    finally:
        if vba:
            vba.close()
    return links


def extract_links_from_pdf(filepath):
    """Extracts links from annotations and visible text in a PDF file."""
    links = []
    try:
        doc = fitz.open(filepath)
        for page in doc:
            for link in page.get_links():
                uri = link.get("uri", "")
                if uri and "http" in uri:
                    links.append(uri.strip())
            text = page.get_text("text")
            found_urls = re.findall(r"https?://\S+|www\.\S+", text)
            for url in found_urls:
                cleaned_url = url.rstrip(".,;)!?\"'")
                links.append(cleaned_url)
        doc.close()
    except Exception as e:
        print(f"[WARNING] Error reading PDF file {filepath}: {e}")
    return links


def extract_links_from_pptx(filepath):
    """Extracts plain text URLs from a pptx file"""

    links = []
    try:
        prs = pptx.Presentation(filepath)
        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                # Check textual content
                if hasattr(shape, "text"):
                    found_urls = re.findall(r"https?://\S+|www\.\S+", shape.text)
                    for url in found_urls:
                        cleaned = url.rstrip(".,;)!?\"'")
                        print(f"[PPTX-TEXT] Slide {slide_num}: {cleaned}")
                        links.append(cleaned)

                # Check for click-action hyperlinks
                if (
                    hasattr(shape, "click_action")
                    and shape.click_action.hyperlink
                    and shape.click_action.hyperlink.address
                ):
                    address = shape.click_action.hyperlink.address.strip()
                    print(f"[PPTX-LINK] Slide {slide_num}: {address}")
                    links.append(address)
    except Exception as e:
        print(f"[WARNING] Error reading PPTX file {filepath}: {e}")
    return links


def extract_links_from_ppt(filepath):
    links = []
    try:
        result = subprocess.run(["strings", filepath], capture_output=True, text=True)
        if result.returncode == 0:
            for line in result.stdout.splitlines():
                found_urls = re.findall(r"https?://\S+|www\.\S+", line)
                for url in found_urls:
                    cleaned = url.rstrip(".,;)!?\"'")
                    if "schemas.openxmlformats.org" in cleaned:
                        continue  # ❌ Skip known false positives
                    print(f"[PPT] Found: {cleaned}")
                    links.append(cleaned)
    except FileNotFoundError:
        print(
            "[ERROR] 'strings' command not found. Install with: sudo apt install binutils"
        )
    except Exception as e:
        print(f"[WARNING] Error reading PPT file {filepath}: {e}")
    return links


def extract_links_from_xlsx(filepath):
    """Extracts visible hyperlinks from cell values in an .xlsx spreadsheet."""

    links = []
    try:
        result = subprocess.run(["strings", filepath], capture_output=True, text=True)
        if result.returncode == 0:
            for line in result.stdout.splitlines():
                found_urls = re.findall(r"https?://\S+|www\.\S+", line)
                for url in found_urls:
                    cleaned = url.rstrip(".,;)!?\"'")
                    if "schemas.openxmlformats.org" in cleaned:
                        continue  # ❌ Skip known false positives
                    print(f"[PPT] Found: {cleaned}")
                    links.append(cleaned)
    except FileNotFoundError:
        print(
            "[ERROR] 'strings' command not found. Install with: sudo apt install binutils"
        )
    except Exception as e:
        print(f"[WARNING] Error reading PPT file {filepath}: {e}")
    return links


def extract_links_from_xlsx(filepath):
    links = []
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        found_urls = re.findall(r"https?://\S+|www\.\S+", cell.value)
                        for url in found_urls:
                            links.append(url.rstrip(".,;)!?\"'"))
    except Exception as e:
        print(f"[WARNING] Error reading XLSX file {filepath}: {e}")
    return links


def extract_links_from_csv(filepath):
    """Extracts URLs from each cell in a CSV file."""

    links = []
    try:
        with open(filepath, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                for cell in row:
                    found_urls = re.findall(r"https?://\S+|www\.\S+", cell)
                    for url in found_urls:
                        links.append(url.rstrip(".,;)!?\"'"))
    except Exception as e:
        print(f"[WARNING] Error reading CSV file {filepath}: {e}")
    return links


def extract_zip_recursive(zip_path, parent_display_path, all_links):
    """Recursively extracts a ZIP file and scans supported files within for links."""

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            with zipfile.ZipFile(zip_path, "r") as z:
                z.extractall(temp_dir)

            for root, dirs, files in os.walk(temp_dir):
                for fname in files:
                    fpath = os.path.join(root, fname)
                    ext = os.path.splitext(fpath)[1].lower()
                    display_path = os.path.relpath(fpath, temp_dir)
                    display_path = f"{parent_display_path}/{display_path}"
                    if ext in SUPPORTED_EXTENSIONS:
                        if ext == ".zip":
                            extract_zip_recursive(fpath, display_path, all_links)
                        else:
                            file_links = scan_single_file(fpath)
                            all_links[display_path] = file_links
    except Exception as e:
        print(f"[WARNING] Failed to handle ZIP {zip_path}: {e}")


def scan_single_file(filepath):
    """Choose link extraction based on file extension."""

    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".docx":
        return extract_links_from_docx(filepath)
    elif ext == ".doc":
        return extract_links_from_doc(filepath)
    elif ext == ".pdf":
        return extract_links_from_pdf(filepath)
    elif ext == ".pptx":
        return extract_links_from_pptx(filepath)
    elif ext == ".ppt":
        return extract_links_from_ppt(filepath)
    elif ext == ".xlsx":
        return extract_links_from_xlsx(filepath)
    elif ext == ".csv":
        return extract_links_from_csv(filepath)
    return []


def scan_folder_for_links(folder_path):
    """Recursively scans all supported files in a folder and extracts links."""

    all_links = {}

    def recursive_scan(path):
        if os.path.isdir(path):
            for item in os.listdir(path):
                full_path = os.path.join(path, item)
                recursive_scan(full_path)
        elif os.path.isfile(path):
            ext = os.path.splitext(path)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                display_name = os.path.relpath(path, folder_path)
                if ext == ".zip":
                    extract_zip_recursive(path, display_name, all_links)
                else:
                    links = scan_single_file(path)
                    all_links[display_name] = links

    recursive_scan(folder_path)
    return all_links


# ---------- Unzip Initial Files ----------
def unzipFolders():
    """Unzips all .zip files in the toProcessFurther folder for further analysis."""
    for filename in os.listdir(TO_PROCESS_FURTHER_PATH):
        if filename.endswith(".zip"):
            zip_path = os.path.join(TO_PROCESS_FURTHER_PATH, filename)
            extract_folder = os.path.join(TO_PROCESS_FURTHER_PATH, filename[:-4])
            os.makedirs(extract_folder, exist_ok=True)
            with zipfile.ZipFile(zip_path, "r") as zip_ref:
                zip_ref.extractall(extract_folder)
            print(f"[INFO] Unzipped: {filename} -> {extract_folder}")


# ---------- ClamAV Scanner ----------
def scanWithClamAV():
    """Runs a ClamAV antivirus scan on the toProcessFurther folder."""
    print("\n[INFO] Running ClamAV scan...")
    scan_command = [
        "clamscan",
        "-r",
        "--infected",
        "--no-summary",
        TO_PROCESS_FURTHER_PATH,
    ]
    try:
        result = subprocess.run(
            scan_command, capture_output=True, text=True, check=False
        )
        print("[INFO] Scan complete.\n")
        print(result.stdout if result.stdout else "[INFO] No infected files found.")
    except FileNotFoundError:
        print("[ERROR] clamscan not found. Install it with: sudo apt install clamav")


def clear_directory():
    """Deletes all files and subfolders in the toProcessFurther directory."""

    if not os.path.exists(TO_PROCESS_FURTHER_PATH):
        print(f"[WARNING] Path does not exist: {TO_PROCESS_FURTHER_PATH}")
        return
    for item in os.listdir(TO_PROCESS_FURTHER_PATH):
        item_path = os.path.join(TO_PROCESS_FURTHER_PATH, item)
        try:
            if os.path.isfile(item_path) or os.path.islink(item_path):
                os.unlink(item_path)
                print(f"[INFO] Deleted file: {item_path}")
            elif os.path.isdir(item_path):
                shutil.rmtree(item_path)
                print(f"[INFO] Deleted folder: {item_path}")
        except Exception as e:
            print(f"[ERROR] Could not delete {item_path}. Reason: {e}")


def normalize_url(url):
    """Normalizes a URL by enforcing scheme and trailing slash rules."""

    if not url.startswith("http://") and not url.startswith("https://"):
        url = "https://" + url

    parsed = urlparse(url)

    # Step 2: Check if path is missing trailing slash
    path = parsed.path

    # Skip if path ends with file-like extension
    if not path.endswith("/") and not re.search(r"\.[a-zA-Z0-9]{2,5}$", path):
        path += "/"

    # Rebuild the full URL
    normalized = urlunparse(
        (
            parsed.scheme,
            parsed.netloc,
            path,
            parsed.params,
            parsed.query,
            parsed.fragment,
        )
    )
    return normalized


# main function
def downloadFilesAndCheck():
    print("=" * 50)
    print("---Processing Downloaded Files---")
    unzipFolders()
    scanWithClamAV()
    found_links = scan_folder_for_links(TO_PROCESS_FURTHER_PATH)

    print("\n[INFO] Scan Summary:")
    all_scanned_files = set()
    all_links = []

    for root, dirs, files in os.walk(TO_PROCESS_FURTHER_PATH):
        for fname in files:
            ext = os.path.splitext(fname)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                rel_path = os.path.relpath(
                    os.path.join(root, fname), TO_PROCESS_FURTHER_PATH
                )
                all_scanned_files.add(rel_path)

    for filepath in sorted(all_scanned_files):
        links = found_links.get(filepath, [])
        print(f"\n[FILE] {filepath} - Found {len(links)} link(s):")
        for link in links:
            normalized = normalize_url(link)
            print(f"   [LINK] {normalized}")
            all_links.append(normalized)

    print("\n[INFO] Total Links Found:")
    for link in all_links:
        print(f"   [SCRAPE] {link}")

    clear_directory()
    print("=" * 50)
    return all_links


if __name__ == "__main__":
    downloadFilesAndCheck()
