from urllib.parse import urlparse
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import re
import asyncio
import os
from urllib.parse import urljoin, urlparse, parse_qs, urlencode, urlunparse
from dotenv import load_dotenv
from playwright.async_api import async_playwright, Page
import requests
from urllib.parse import urljoin, urlparse
import aiohttp
import filetype  # pip install filetype
from datetime import datetime, UTC

# --- Load credentials ---
load_dotenv()
USERNAME = os.getenv("MOODLE_USERNAME")
PASSWORD = os.getenv("MOODLE_PASSWORD")
print(os.path)

print(PASSWORD)

MOODLE_DOMAIN = "3.107.195.248"

EXCLUDED_PATH_PREFIXES = [
    "/moodle/user/",
    "/moodle/message/",
    "/moodle/notes/",
    "/moodle/blog/",
    "/moodle/iplookup/",
    "/moodle/tag/",
    "/moodle/calendar/",
    "/moodle/report/usersessions/",
    "/moodle/admin/",
    "/moodle/enrol/",
    "/moodle/grade/report/overview/",
    "/moodle/competency/",
]

def should_exclude_url(url: str) -> bool:
    path = urlparse(url).path
    return any(path.startswith(prefix) for prefix in EXCLUDED_PATH_PREFIXES)

def is_internal(url: str) -> bool:
    return urlparse(url).netloc == MOODLE_DOMAIN

def is_external(url: str) -> bool:
    return not is_internal(url) or should_exclude_url(url)

def extract_links_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        all_links = []

        for page in reader.pages:
            if "/Annots" in page:
                for annot in page["/Annots"]:
                    obj = annot.get_object()
                    uri = obj.get("/A", {}).get("/URI")
                    if uri:
                        all_links.append(uri)

        print(f"🔗 All Links Found in {pdf_path}:")
        for link in all_links:
            print(" -", link)

        external_links = [link for link in all_links if is_external(link)]

        print("\n🌐 External Links (filtered by is_external):")
        for link in external_links:
            print(" -", link)

        return all_links, external_links

    except Exception as e:
        print(f"❌ Failed to extract links from {pdf_path}: {e}")
        return [], []
    


URL_REGEX = re.compile(r"https?://[^\s]+")

def extract_links_from_docx(docx_path):
    doc = Document(docx_path)
    all_links = set()

    # Method 1: Hyperlinks in relationships
    for rel in doc.part.rels.values():
        if "hyperlink" in rel.reltype:
            all_links.add(rel.target_ref)

    # Method 2: Regex-based scan of visible text
    for para in doc.paragraphs:
        found = URL_REGEX.findall(para.text)
        all_links.update(found)

    external_links = [link for link in all_links if is_external(link)]

    print("🔗 All Links Found in .docx:")
    for link in all_links:
        print(" -", link)

    print("\n🌐 External Links:")
    for link in external_links:
        print(" -", link)

    return list(all_links), external_links

def extract_links_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    all_links = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            # Method 1: Check for shape with hyperlink
            if shape.has_text_frame and shape.click_action.hyperlink.address:
                all_links.add(shape.click_action.hyperlink.address)

            # Method 2: Scan text for URLs
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        found = URL_REGEX.findall(run.text)
                        all_links.update(found)

    external_links = [link for link in all_links if is_external(link)]

    print("🔗 All Links Found in .pptx:")
    for link in all_links:
        print(" -", link)

    print("\n🌐 External Links:")
    for link in external_links:
        print(" -", link)

    return list(all_links), external_links

extract_links_from_pdf("temp/tempFile.pdf")
extract_links_from_docx("temp/tempFile.docx")
extract_links_from_pptx("temp/tempFile.pptx")